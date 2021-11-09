# powerpoint に colormap, 編集済みpic, title, 計算したYe などを出力するプログラム
print("run 'auto_pptx_pic.py'")

# ------------------------------------------------------------------------------------------------------------------------------
from matplotlib.pyplot import colormaps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
import os
from posixpath import normcase
from config import *

# -----------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
# PATH の設定
PIC_EDIT_PATH = PATH + "pic_edit/"
COLORMAP_PATH = PATH + "result/colormap/"
FIG_A_MF_PATH = PATH + "result/FIG_A_MF/"
DATA_PATH = PATH + "data/"
PPTX_PATH = PATH + "result/pic_pptx/"
Ye_FILE = PATH + "result/Ye_sumMF.txt"

# -----------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
f = open(Ye_FILE, 'r')
DATALIST = f.readlines()

if not os.path.exists(PPTX_PATH):
    os.makedirs(PPTX_PATH)

FILE2 = PATH + "result/cnt_larger_100.txt"
# if os.path.exists(FILE2):
#     os.remove(FILE2)

# -----------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
# 実行
f2 = open(FILE2, "r")
for w in status:

    for x in Ye:
        n = 0
        # -----------------------------------------------------------------------
        # Presentation オブジェクトを生成
        ppt = Presentation()

        # スライドのサイズを指定
        ppt.slide_width = Cm(33.868)
        ppt.slide_heght = Cm(19.05)

        # ----------------------------------------------------------------------
        # スライド1枚目
        # 追加するスライドを選択
        slide_layout_0 = ppt.slide_layouts[0]
        # スライドを追加
        slide_0 = ppt.slides.add_slide(slide_layout_0)
        # place holder (テキスト入れるやつ？)を選択
        slide_0_title = slide_0.placeholders[0]
        # スライドに貼り付け
        print(type(slide_0_title))
        # スライドの幅
        slide_0_title.width = Cm(30)
        # スライドの高さ
        slide_0_title.height = Cm(20)
        # スライドの文字
        slide_0_title.text = "result picture"

        # ----------------------------------------------------------------------
        # スライド2枚目
        # 追加するスライドを選択
        slide_layout_0 = ppt.slide_layouts[0]
        # スライドを追加
        slide_0 = ppt.slides.add_slide(slide_layout_0)
        # place holder (テキスト入れるやつ？)を選択
        slide_0_title = slide_0.placeholders[0]
        # スライドに貼り付け
        print(type(slide_0_title))
        # スライドの幅
        slide_0_title.width = Cm(30)
        # スライドの高さ
        slide_0_title.height = Cm(20)
        # スライドの文字
        slide_0_title.text = "Ye:" + x
        # ----------------------------------------------------------------------
        for y in T0:
            m = 0
            for z in rho0:
                # ----------------------------------------------------------------------
                # スライド3枚目以降 (それぞれの条件でのcolormapとスクショ)

                # TITLE, PATHの指定
                TI = "Ye_" + x + "_T0_" + y + "_rho0_" + z
                TITLE = w + "_" + TI
                INPUT_PIC_PATH = PIC_EDIT_PATH + TITLE + ".png"
                INPUT_CLMP_PATH = COLORMAP_PATH + "Thorium232/last_Thorium232_Ye_" + x + ".png"
                INPUT_A_MF_PATH = FIG_A_MF_PATH + TITLE + ".png"
                
                if not os.path.exists(INPUT_PIC_PATH):
                    m += 1
                    continue

                # ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
                # ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
                # スライドのレイアウトを選択
                slide_layout_6 = ppt.slide_layouts[6]
                # スライドを追加
                slide_6 = ppt.slides.add_slide(slide_layout_6)

                # 画像のplaceholder を作成
                shapes = slide_6.shapes

                shapes.add_picture(INPUT_PIC_PATH, Cm(15), Cm(2), width = None, height = Cm(8))
                shapes.add_picture(INPUT_CLMP_PATH, 0, Cm(3), width = None, height = Cm(10.88))
                shapes.add_picture(INPUT_A_MF_PATH, Cm(15), Cm(10), width = None, height = Cm(8))
                
                # 指定する長方形を追加
                rect0 = slide_6.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Pt(52 + 23 * m), Pt(122 + 26 * n),
                    Pt(22), Pt(26)
                )
                rect0.fill.background()
                rect0.line.width = Pt(2)

                # # placeholder_1 = slide_6.placeholders[1]
                shape = shapes.add_textbox(Cm(1), Cm(1), Cm(22), Cm(16))
                # # box内にテキストを追加
                shape.text = INPUT_PIC_PATH

                # ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
                # Ye_sum.txt から当てはまる条件のYeとMFを抽出。
                NOW_sumYe = ""
                NOW_sumMF = ""
                for data in DATALIST:
                    status_, Ye_, T0_, rho0_, sumYe, sumMF = data.split()
                    if w == status_ and x == Ye_ and y == T0_ and z == rho0_:
                        NOW_sumYe = sumYe
                        NOW_sumMF = sumMF
                        
                # ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
                # Ye, sum_MF を入力するtextbox を作成
                shape1 = shapes.add_textbox(Cm(2), Cm(14),Cm(4), Cm(15))
                shape1.text = "Ye : " + NOW_sumYe
                
                shape2 = shapes.add_textbox(Cm(2), Cm(15),Cm(4), Cm(16))
                shape2.text = "sum_MF : " + NOW_sumMF
                
                # ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
                # Informationdataの下から数行を読み取るプログラムを書く。
                cnt = 0
                cntMAX = 0
                if w == "last":
                    for i in range(101):
                        INFORMATION_PATH = DATA_PATH + "InformationData_" + TI + "_cnt_" + str(cnt) + ".txt"
                        if os.path.exists(INFORMATION_PATH):
                            cntMAX = max(int(cntMAX), int(cnt))
                        cnt += 1
                    # if cntMAX >= 100:
                    #     print("cnt100 : %s" % TITLE)
                    #     f2.write("%s\n" % TITLE)
                shape3 = shapes.add_textbox(Cm(2), Cm(16),Cm(4), Cm(17))
                shape3.text = "[Information] cnt :  " + str(cntMAX)
                m += 1

                # ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
            n += 1
            # スライドのレイアウトを選択
            slide_layout_6 = ppt.slide_layouts[6]
            # スライドを追加
            slide_6 = ppt.slides.add_slide(slide_layout_6)

        # 保存  
        ppt.save(PPTX_PATH + w + "_Ye_"+ x + ".pptx")